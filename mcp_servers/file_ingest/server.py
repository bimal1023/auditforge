"""
File Ingest MCP Server
======================
Exposes parsers for the document types a PE data room throws at us:

  parse_pdf(file_path)   → extracted text + page count
  parse_csv(file_path)   → rows as list of dicts + row count
  parse_xlsx(file_path)  → all sheets rendered to tab-separated text
  parse_docx(file_path)  → paragraphs + table text
  parse_pptx(file_path)  → per-slide shape text + tables + speaker notes
  parse_txt(file_path)   → raw UTF-8 text

Run standalone:
  python mcp_servers/file_ingest/server.py
"""
from __future__ import annotations

import csv
from pathlib import Path

from mcp.server.fastmcp import FastMCP

mcp = FastMCP("file-ingest")

MAX_TEXT_CHARS = 100_000   # cap extracted text to avoid flooding context / DB

# Only allow reading from OS temp directories — blocks path traversal attempts
import tempfile as _tempfile
# Resolve symlinks for each candidate dir (macOS: /tmp → /private/tmp, /var → /private/var)
_ALLOWED_DIRS = frozenset(
    str(Path(d).resolve())
    for d in {_tempfile.gettempdir(), "/tmp", "/var/folders", "/var/tmp"}
    if Path(d).exists()
)


def _assert_safe_path(file_path: str) -> Path:
    path = Path(file_path).resolve()
    if not any(str(path).startswith(d) for d in _ALLOWED_DIRS):
        raise ValueError(f"Access denied: path must be inside a temp directory")
    return path


@mcp.tool()
async def parse_pdf(file_path: str) -> dict:
    """
    Extract text from a PDF file using pypdf.

    Args:
        file_path : absolute or relative path to the .pdf file

    Returns:
        text      : concatenated text from all pages
        pages     : total page count
        file      : echoed file path
    """
    try:
        import pypdf  # lazy import so the server starts even without pypdf
    except ImportError:
        raise RuntimeError("pypdf is not installed — run: pip install pypdf")

    try:
        path = _assert_safe_path(file_path)
    except ValueError as e:
        return {"error": str(e)}
    if not path.exists():
        return {"error": "File not found"}
    if path.suffix.lower() != ".pdf":
        return {"error": "Not a PDF file"}

    reader = pypdf.PdfReader(str(path))
    pages_text: list[str] = []
    for page in reader.pages:
        extracted = page.extract_text()
        if extracted:
            pages_text.append(extracted)

    full_text = "\n\n".join(pages_text)
    if len(full_text) > MAX_TEXT_CHARS:
        full_text = full_text[:MAX_TEXT_CHARS]

    return {
        "text": full_text,
        "pages": len(reader.pages),
        "file": str(path),
    }


@mcp.tool()
async def parse_csv(file_path: str) -> dict:
    """
    Parse a CSV file using the stdlib csv module.

    Args:
        file_path : absolute or relative path to the .csv file

    Returns:
        headers   : list of column names
        rows      : list of row dicts (column → value)
        count     : total row count
        file      : echoed file path
    """
    try:
        path = _assert_safe_path(file_path)
    except ValueError as e:
        return {"error": str(e)}
    if not path.exists():
        return {"error": "File not found"}

    rows: list[dict] = []
    headers: list[str] = []

    with open(path, newline="", encoding="utf-8-sig") as fh:
        reader = csv.DictReader(fh)
        headers = list(reader.fieldnames or [])
        for row in reader:
            rows.append(dict(row))

    return {
        "headers": headers,
        "rows": rows,
        "count": len(rows),
        "file": str(path),
    }


@mcp.tool()
async def parse_xlsx(file_path: str) -> dict:
    """
    Extract text from an Excel workbook (.xlsx) using openpyxl.

    Every sheet is rendered as a "## Sheet: <name>" heading followed by its
    rows joined with tabs — a flat, RAG-friendly text representation that keeps
    the numbers and their column context together.

    Args:
        file_path : path to the .xlsx file

    Returns:
        text      : concatenated text across all sheets (capped)
        sheets    : number of sheets
        rows      : number of non-empty rows rendered
        file      : echoed file path
    """
    try:
        import openpyxl  # lazy import so the server starts even without openpyxl
    except ImportError:
        raise RuntimeError("openpyxl is not installed — run: pip install openpyxl")

    try:
        path = _assert_safe_path(file_path)
    except ValueError as e:
        return {"error": str(e)}
    if not path.exists():
        return {"error": "File not found"}
    if path.suffix.lower() != ".xlsx":
        return {"error": "Not an XLSX file"}

    try:
        wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    except Exception as e:  # noqa: BLE001 — surface any openpyxl error to the caller
        return {"error": f"Could not open workbook: {e}"}

    sheet_count = len(wb.sheetnames)
    parts: list[str] = []
    total_rows = 0
    running_len = 0
    truncated = False
    for ws in wb.worksheets:
        header = f"## Sheet: {ws.title}"
        parts.append(header)
        running_len += len(header) + 1
        for row in ws.iter_rows(values_only=True):
            cells = ["" if v is None else str(v) for v in row]
            if not any(c.strip() for c in cells):
                continue
            line = "\t".join(cells)
            if running_len + len(line) + 1 > MAX_TEXT_CHARS:
                truncated = True
                break
            parts.append(line)
            running_len += len(line) + 1
            total_rows += 1
        if truncated:
            break
    wb.close()

    text = "\n".join(parts)[:MAX_TEXT_CHARS]
    return {
        "text": text,
        "sheets": sheet_count,
        "rows": total_rows,
        "file": str(path),
    }


@mcp.tool()
async def parse_docx(file_path: str) -> dict:
    """
    Extract text from a Word document (.docx) using python-docx.

    Pulls body paragraphs in order, then appends any table rows (tab-joined)
    so figures embedded in tables aren't lost.

    Args:
        file_path : path to the .docx file

    Returns:
        text       : concatenated paragraph + table text (capped)
        paragraphs : number of non-empty paragraphs
        file       : echoed file path
    """
    try:
        import docx  # python-docx; lazy import
    except ImportError:
        raise RuntimeError("python-docx is not installed — run: pip install python-docx")

    try:
        path = _assert_safe_path(file_path)
    except ValueError as e:
        return {"error": str(e)}
    if not path.exists():
        return {"error": "File not found"}
    if path.suffix.lower() != ".docx":
        return {"error": "Not a DOCX file"}

    try:
        document = docx.Document(str(path))
    except Exception as e:  # noqa: BLE001
        return {"error": f"Could not open document: {e}"}

    parts: list[str] = []
    para_count = 0
    for para in document.paragraphs:
        t = para.text.strip()
        if t:
            parts.append(t)
            para_count += 1

    for table in document.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells]
            if any(cells):
                parts.append("\t".join(cells))

    text = "\n".join(parts)[:MAX_TEXT_CHARS]
    return {
        "text": text,
        "paragraphs": para_count,
        "file": str(path),
    }


@mcp.tool()
async def parse_pptx(file_path: str) -> dict:
    """
    Extract text from a PowerPoint deck (.pptx) using python-pptx.

    Walks every slide in order, pulling text from all shapes (titles, body
    placeholders, text boxes), table cells (tab-joined), and the slide's
    speaker notes — the narrative an analyst most wants the agents to cite.
    Each slide is prefixed with a "## Slide N" heading.

    Args:
        file_path : path to the .pptx file

    Returns:
        text   : concatenated per-slide text (capped)
        slides : number of slides
        file   : echoed file path
    """
    try:
        from pptx import Presentation  # python-pptx; lazy import
    except ImportError:
        raise RuntimeError("python-pptx is not installed — run: pip install python-pptx")

    try:
        path = _assert_safe_path(file_path)
    except ValueError as e:
        return {"error": str(e)}
    if not path.exists():
        return {"error": "File not found"}
    if path.suffix.lower() != ".pptx":
        return {"error": "Not a PPTX file"}

    try:
        prs = Presentation(str(path))
    except Exception as e:  # noqa: BLE001
        return {"error": f"Could not open presentation: {e}"}

    parts: list[str] = []
    running_len = 0
    slide_count = 0
    truncated = False

    def _add(line: str) -> bool:
        """Append a line, respecting MAX_TEXT_CHARS. Returns False once full."""
        nonlocal running_len
        if not line:
            return True
        if running_len + len(line) + 1 > MAX_TEXT_CHARS:
            return False
        parts.append(line)
        running_len += len(line) + 1
        return True

    for idx, slide in enumerate(prs.slides, start=1):
        slide_count += 1
        if not _add(f"## Slide {idx}"):
            truncated = True
            break
        for shape in slide.shapes:
            # Tables: render each row tab-joined.
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells]
                    if any(cells) and not _add("\t".join(cells)):
                        truncated = True
                        break
                if truncated:
                    break
                continue
            # Text frames: titles, body placeholders, text boxes.
            if shape.has_text_frame:
                t = shape.text_frame.text.strip()
                if t and not _add(t):
                    truncated = True
                    break
        if truncated:
            break
        # Speaker notes — often the richest narrative.
        if slide.has_notes_slide:
            notes = (slide.notes_slide.notes_text_frame.text or "").strip()
            if notes and not _add(f"[Notes] {notes}"):
                truncated = True
        if truncated:
            break

    text = "\n".join(parts)[:MAX_TEXT_CHARS]
    return {
        "text": text,
        "slides": slide_count,
        "file": str(path),
    }


@mcp.tool()
async def parse_txt(file_path: str) -> dict:
    """
    Read a plain-text file (.txt) as UTF-8.

    Args:
        file_path : path to the .txt file

    Returns:
        text  : file contents (capped)
        chars : total character count of the original file
        file  : echoed file path
    """
    try:
        path = _assert_safe_path(file_path)
    except ValueError as e:
        return {"error": str(e)}
    if not path.exists():
        return {"error": "File not found"}

    try:
        raw = path.read_text(encoding="utf-8", errors="replace")
    except Exception as e:  # noqa: BLE001
        return {"error": f"Could not read file: {e}"}

    return {
        "text": raw[:MAX_TEXT_CHARS],
        "chars": len(raw),
        "file": str(path),
    }


if __name__ == "__main__":
    mcp.run(transport="stdio")
